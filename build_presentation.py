#!/usr/bin/env python3
"""Build PhD oral-examination presentation from sample template + candidate content."""

from __future__ import annotations

import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "Presentation-sample_EN.pptx"
OUTPUT = ROOT / "PhD-Presentation_EN.pptx"
IMG_DIR = ROOT / "motivation_images"

TITLE = (
    "Doctoral Program Oral Examination\n\n"
    "– Deployable Hybrid Quantum–Classical Edge Intelligence "
    "for Sustainable IoT Security –"
)
SUBTITLE = "Sajib Kumar Mitra"


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def set_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text


def normalize_ws(text: str) -> str:
    return re.sub(r"[\x0b\r\n]+", " ", text).strip()


def replace_in_shape(shape, mapping: dict[str, str]) -> None:
    if not shape.has_text_frame:
        return
    text = shape.text_frame.text
    if not text.strip():
        return
    new_text = text
    for old, new in mapping.items():
        if old in new_text:
            new_text = new_text.replace(old, new)
    if new_text != text:
        set_text(shape, new_text)


def replace_shape_if_contains(shape, needles: tuple[str, ...], new_text: str) -> bool:
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text
    norm = normalize_ws(text)
    if any(needle in norm or needle in text for needle in needles):
        set_text(shape, new_text)
        return True
    return False


def replace_all_text(slide, mapping: dict[str, str]) -> None:
    for shape in iter_shapes(slide.shapes):
        replace_in_shape(shape, mapping)


def replace_title_content(slide, title: str | None = None, body: str | None = None) -> None:
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        idx = shape.placeholder_format.idx
        if idx == 0 and title is not None:
            set_text(shape, title)
        elif idx == 1 and body is not None:
            set_text(shape, body)


def add_image(slide, path: Path, left, top, width) -> None:
    slide.shapes.add_picture(str(path), left, top, width=width)


def clear_slide_pictures(slide) -> None:
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            sp_tree.remove(shape._element)


def simplify_slide(slide, title: str, body: str, image: Path | None = None) -> None:
    """Keep only title/number placeholders; replace with clean title, body, optional image."""
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            if idx == 0:
                set_text(shape, title)
            elif idx == 1:
                set_text(shape, "")
            continue
        sp_tree.remove(shape._element)

    top = Inches(1.45)
    if image and image.exists():
        add_image(slide, image, Inches(0.7), top, Inches(5.8))
        text_left = Inches(6.8)
        text_width = Inches(5.5)
    else:
        text_left = Inches(0.7)
        text_width = Inches(11.6)

    if body:
        box = slide.shapes.add_textbox(text_left, top, text_width, Inches(5.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = body
        p.font.size = Pt(17)


def fix_template_remnants(prs: Presentation) -> None:
    """Remove Confidential Computing sample text that uses vertical-tab line breaks."""
    slide11_fixes = [
        (
            ("Trusted Execution Environment", "TEE"),
            "Billions of IoT endpoints generate continuous telemetry under strict "
            "energy, latency, and trust constraints.",
        ),
        (
            ("Confidential Computing Consortium", "CCC"),
            "QML toolchains (PennyLane, Qiskit) support hybrid quantum–classical pipelines.",
        ),
        (
            ("Confidential Computing ( CC )", "Data in Use"),
            "Hybrid Quantum–Classical Edge Intelligence: train on simulator/cloud, infer on edge.",
        ),
        (
            ("advent of the digital age", "data exchange"),
            "Edge AI reduces backhaul and improves responsiveness; adversarial analytics "
            "motivate selective Quantum AI in hybrid pipelines.",
        ),
        (
            ("privacy and security", "is essential"),
            "Sustainable IoT security analytics must co-design accuracy, circuit depth, and deployability.",
        ),
    ]
    for needles, new in slide11_fixes:
        for shape in iter_shapes(prs.slides[10].shapes):
            replace_shape_if_contains(shape, needles, new)

    slide12_fixes = [
        (("TEE", "Trusted Execution Environment"), "Hybrid QML\n(Edge Intelligence)"),
        (("CC", "Confidential Computing"), "IoT Security\n(Anomaly Detection)"),
    ]
    for needles, new in slide12_fixes:
        for shape in iter_shapes(prs.slides[11].shapes):
            replace_shape_if_contains(shape, needles, new)
    for shape in iter_shapes(prs.slides[11].shapes):
        replace_shape_if_contains(
            shape,
            ("Remote Attestation", "Benchmark harness verification"),
            "Local anomaly scoring with compressed quantum–classical models",
        )

    slide13_fixes = [
        (
            ("To enhance CC", "Remote Attestation"),
            "To deliver deployable hybrid edge intelligence, we address four research gaps in IoT QML:",
        ),
        (
            ("the TEE usage environment", "Chain of If"),
            "Gap 1: Cloud-centric training — VQCs sized for QPUs, not microcontroller/FPGA budgets.",
        ),
        (
            ("Advanced CC", "orchestrator integration"),
            "Gaps 2–4: Missing metric alignment, weak deployment path, and over-broad quantum claims.",
        ),
        (
            ("complex CoT", "extending RA"),
            "Doctoral aim: sustainable IoT anomaly detection under fixed power and latency budgets.",
        ),
    ]
    for needles, new in slide13_fixes:
        for shape in iter_shapes(prs.slides[12].shapes):
            replace_shape_if_contains(shape, needles, new)

    slide15_fixes = [
        (
            ("Compressed VQC inference", "TEE"),
            "Compression &\nexport",
        ),
        (
            ("CC complexly integrated", "Remote Attestation"),
            "Two-phase pipeline: train on simulator/cloud, compress, and infer on edge.",
        ),
        (
            ("Isolation Execution Environment",),
            "Hybrid Edge QML",
        ),
        (
            ("advancement of CC",),
            "Empirical guidance on when hybrid Quantum AI outperforms classical edge AI.",
        ),
    ]
    for needles, new in slide15_fixes:
        for shape in iter_shapes(prs.slides[14].shapes):
            replace_shape_if_contains(shape, needles, new)

    # Slide 16: remove AMD attestation reference captions
    for shape in iter_shapes(prs.slides[15].shapes):
        if shape.has_text_frame:
            t = normalize_ws(shape.text_frame.text)
            if any(k in t for k in ("AMD SEV", "amd.com", "Attestation: Establishing")):
                set_text(shape, "")

    # Slide 17: replace CC stakeholder boilerplate
    s17_fixes = [
        (
            ("TEE hardware vendors", "CC cloud providers"),
            "Stakeholders: supervisor, FPGA lab, IoT testbed operators, open-source QML community",
        ),
        (
            ("distributed computing resources", "GPU -based TEE"),
            "Use cases: smart-building gateway IDS, industrial IoT anomaly detection, post-discharge RPM IoMT",
        ),
        (
            ("conventional TEEs and CCs", "RATs"),
            "Investigate prior QML for IoT, projected quantum kernels, and survey benchmark gaps",
        ),
        (
            ("extension method for CoT", "Remote Attestation"),
            "Propose compressed VQC / quantum-kernel models with explicit qubit-depth budgets",
        ),
        (
            ("cryptographic and key management", "RATs"),
            "Design compression, export, and edge inference flows using PennyLane/Qiskit toolchains",
        ),
        (
            ("PKI", "TLS"),
            "Integrate with public IoT datasets, power metering, and reproducible benchmark harness",
        ),
    ]
    for needles, new in s17_fixes:
        for shape in iter_shapes(prs.slides[16].shapes):
            replace_shape_if_contains(shape, needles, new)

    # Slide 14 callout bubbles
    for shape in iter_shapes(prs.slides[13].shapes):
        replace_shape_if_contains(
            shape,
            ("the TEE usage environment",),
            "Survey gap:\nCloud-centric QML",
        )
        replace_shape_if_contains(
            shape,
            ("Advanced CC", "Gaps 2"),
            "Doctoral focus:\nEdge deployability",
        )

    # Slide 19 timeline boilerplate
    for shape in iter_shapes(prs.slides[18].shapes):
        replace_shape_if_contains(
            shape,
            ("examining technical requirements", "doctoral dissertation"),
            "Iterative loop: benchmark → hybrid model → FPGA prototype → thesis",
        )

    # Slide 21 divider (section slide — title placeholder absent in template)
    s21 = prs.slides[20]
    for shape in list(s21.shapes):
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            set_text(
                shape,
                "Supplementary slides follow\n\n"
                "Survey highlights · Four paradigms · Benchmark alignment · TRL roadmap",
            )
        elif not shape.is_placeholder:
            s21.shapes._spTree.remove(shape._element)


def build() -> Path:
    shutil.copy2(TEMPLATE, OUTPUT)
    prs = Presentation(str(OUTPUT))

    # --- Slide 1: Title ---
    s1 = prs.slides[0]
    replace_title_content(
        s1,
        title=TITLE,
        body=SUBTITLE,
    )

    # --- Slide 2: Agenda ---
    replace_title_content(
        prs.slides[1],
        body=(
            "Previous work experience and job responsibilities\n"
            "Research achievements (M.S. thesis, published papers, ongoing survey)\n"
            "Motivation for pursuing doctoral study\n"
            "Research plan"
        ),
    )

    # --- Slide 3: Work experience ---
    replace_title_content(
        prs.slides[2],
        body=(
            "B.Sc. in Computer Science & Engineering, University of Dhaka\n"
            "M.S. in Computer Science & Engineering, University of Dhaka\n"
            "　　(Thesis: efficient reversible logic & quantum circuit design)\n\n"
            "Joint Director (ICT), Bangladesh Bank\n"
            "　　(Large-scale secure ICT systems, software engineering, audit support)\n\n"
            "Full-time availability for funded 3–4 year doctoral study"
        ),
    )

    # --- Slide 4: Business content ---
    replace_title_content(prs.slides[3], body="ICT systems engineering & security R&D")

    # --- Slide 5: Responsibilities ---
    replace_title_content(
        prs.slides[4],
        body=(
            "Secure ICT architecture and software engineering leadership\n"
            "Research on quantum computing, reversible logic, and IoT security analytics\n"
            "Cross-functional coordination for national-scale financial ICT systems"
        ),
    )

    # --- Slide 6: Research achievements table ---
    s6 = prs.slides[5]
    replace_title_content(
        s6,
        body=(
            "Research spans quantum computing, reversible logic, hybrid IoT security, "
            "and embedded AI—with publications in IEEE/ACM venues and an ongoing "
            "systematic IoT–Quantum AI survey (co-investigator)."
        ),
    )
    for shape in s6.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        rows = [
            [
                "Reversible / quantum logic design",
                "—",
                "IEEE conf.",
                "3+",
                "—",
            ],
            [
                "Quantum circuits & optimization",
                "—",
                "IEEE / ACM",
                "2+",
                "—",
            ],
            [
                "IoT–Quantum AI survey (ongoing)",
                "Survey",
                "—",
                "1",
                "—",
            ],
            [
                "Hybrid edge QML (doctoral target)",
                "Target",
                "IEEE IoT J., FGCS",
                "—",
                "—",
            ],
            [
                "total (selected)",
                "1",
                "5+",
                "—",
                "—",
            ],
        ]
        for r, row_data in enumerate(rows, start=1):
            for c, val in enumerate(row_data):
                table.cell(r, c).text = val
        break
    for shape in s6.shapes:
        if shape.has_text_frame and "domestic publications" in shape.text_frame.text:
            set_text(shape, "* Counts reflect selected IEEE/ACM and survey outputs.")

    # --- Slide 7: Ongoing survey reference ---
    replace_title_content(
        prs.slides[6],
        title="[ Reference ] Ongoing survey — Quantum AI-powered IoT",
        body="",
    )
    s7 = prs.slides[6]
    refs = [
        "( Survey manuscript )",
        "( PRISMA corpus n=57: 42 IoT + 15 enabling refs. )",
        "• Four paradigms: Conventional, AI, Quantum, Quantum-AI IoT\n"
        "• Unified anomaly-detection benchmarks (latency, complexity, energy proxies)\n"
        "• Smart-building gateway IDS & post-discharge RPM IoMT case studies",
        "• Hybrid Reference Architecture with policy-routed deployment paths α–γ",
    ]
    tb_idx = 0
    for shape in iter_shapes(s7.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
            if "SCIS" in shape.text_frame.text or "CSS" in shape.text_frame.text:
                if tb_idx < len(refs):
                    set_text(shape, refs[tb_idx])
                    tb_idx += 1

    # --- Slide 8: Past publications ---
    replace_title_content(
        prs.slides[7],
        title="[ Reference ] Past presentations (M.S. thesis, published papers)",
    )
    s8 = prs.slides[7]
    pub_map = {
        "( IEICE Japanese journal)": "( IEEE / ACM conferences )",
        "(Master's thesis)": "( M.S. thesis, University of Dhaka )",
        "( Presented at domestic symposia such as SITA, SCIS, and ISEC )": (
            "( Ongoing: IoT–Quantum AI survey with PRISMA corpus )"
        ),
    }
    replace_all_text(s8, pub_map)
    for shape in iter_shapes(s8.shapes):
        if shape.has_text_frame:
            t = shape.text_frame.text
            if "decryption key is split" in t:
                set_text(
                    shape,
                    "• Efficient reversible logic & quantum circuit optimization\n"
                    "• Low-power arithmetic datapaths for embedded quantum-classical co-processors",
                )
            if t.strip() == "・":
                set_text(shape, "• Published in IEEE and ACM conference proceedings")

    # --- Slide 9: Motivation ---
    s9 = prs.slides[8]
    replace_title_content(
        s9,
        body=(
            "To lead security & embedded-AI R&D at the intersection of quantum computing and IoT …\n"
            "Doctoral study will solidify expertise in deployable hybrid edge intelligence.\n\n"
            "Cultivation goals:\n"
            "• Research execution ability on hardware-aware QML pipelines\n"
            "• Technical leadership bridging academic research and national-scale ICT practice\n"
            "• Bridge academic survey insights and reproducible edge prototypes"
        ),
    )
    if (IMG_DIR / "05_motivation_composite.png").exists():
        add_image(s9, IMG_DIR / "05_motivation_composite.png", Inches(7.0), Inches(1.3), Inches(5.8))

    # --- Slide 10: Research plan title ---
    replace_title_content(
        prs.slides[9],
        body=(
            "Deployable Hybrid Quantum–Classical Edge Intelligence\n"
            "for Sustainable IoT Security and Anomaly Detection"
        ),
    )

    # --- Slides 11–15: Background / research area (targeted CC→QAI replacements) ---
    cc_to_qai = {
        "Trusted Execution Environment ( TEE ) , which enables a \n"
        '" reliable execution environment " based on hardware, is a technology attracting attention!':
            "Billions of IoT endpoints generate continuous telemetry under strict energy, "
            "latency, and trust constraints.",
        "a Linux Foundation project,\nThe Confidential Computing Consortium ( CCC ) has been established.":
            "QML toolchains (PennyLane, Qiskit) enable hybrid quantum–classical learning pipelines.",
        "With the advent of the digital age and rapid technological innovation, \n"
        "the need for data exchange that transcends the boundaries of individuals and organizations is rapidly increasing.":
            "Edge AI reduces backhaul and improves responsiveness, yet high-dimensional "
            "streaming workloads motivate selective Quantum AI at the edge.",
        "privacy and security \nis essential.":
            "Sustainable IoT security analytics must co-design accuracy, circuit depth, "
            "and deployability.",
        "Expectations for Privacy Enhancing Technologies": "Quantum AI for IoT Security",
        "Industry Trends": "IoT Trends",
        "technology trends": "Research Trends",
        "SEV": "TinyML",
        "SGX, TDX": "VQC / Q-kernel",
        "TrustZone , CCA": "PennyLane",
        "This has evolved into areas such as Confidential Computing ( CC ) , which aims to protect data during processing ( Data in Use ).":
            "Hybrid Quantum–Classical Edge Intelligence targets train-on-cloud/simulator, infer-on-edge pipelines.",
        "Rapid growth in practical applications": "Rapid QML progress for IoT",
        "From 2019": "2020–2026",
        "the Cloud Native Computing Foundation ,\nThe Confidential Container project has launched.":
            "IEEE / ACM IoT & QCE venues show accelerating hybrid QML research.",
        "TEE \n( Trusted Execution Environment)": "Hybrid QML\n( Edge Intelligence )",
        "CC \n(Confidential Computing)": "IoT Security\n( Anomaly Detection )",
        "Isolation execution environment": "Compressed VQC inference",
        "Chain of Trust": "Benchmark harness",
        "Sensitive data processing": "Streaming IoT features",
        "Signatures, encryption, key management, etc.": "Qiskit, PennyLane, FPGA offload",
        "Memory encryption, etc.\nIsolate the environment,\nPrevent attackers from stealing information.":
            "Circuit depth, qubit budget,\nlatency & energy caps",
        "Isolation execution environment and \nChain of Trust verification \n( Remote Attestation )":
            "Local anomaly scoring with\ncompressed quantum–classical models",
        "Cloud Service \n(Azure, AWS, GCP…)": "IoT gateways, FPGA SoC,\nESP32 sensor emulation",
        "confidential\nContainer": "Sensor\nnodes",
        "confidential\nVM": "Edge\ngateway",
        "TEE -compatible GPU": "FPGA accelerator",
        "NVIDIA H100 , etc.": "Zynq / Artix-class",
        "AMD EPYC , etc.": "ARM IoT gateway",
        "Using confidential containers and confidential VMs": "Using hybrid QML pipelines on constrained edge hardware",
        "Cloud service providers": "Edge platform operators",
        "Along with management functions and peripheral functions,\nProvides confidential containers and confidential VMs .":
            "Trains on simulator/cloud; deploys compressed inference on a single embedded platform.",
        "Applications to server systems are developing.": "Edge deployment paths remain the critical gap.",
        "To enhance CC (Chain of Trust), we will extend Chain of Trust and Remote Attestation .":
            "To deliver deployable hybrid edge intelligence, we address four research gaps in IoT QML:",
        "the TEE usage environment. | There are many different methods for each CC , but Chain of If there are deficiencies in Trust , CC 's security cannot be ensured.":
            "Gap 1: Cloud-centric training — VQCs sized for QPUs, not microcontroller/FPGA budgets.",
        "Advanced CC | CC becomes more sophisticated, enabling features such as orchestrator integration and service development , it becomes necessary to expand the Chain of Trust and implement corresponding ":
            "Gaps 2–4: Missing metric alignment, weak deployment path, and over-broad quantum advantage claims.",
        "Task 1": "Gap 1",
        "Task 2": "Gaps 2–4",
        "of \nthe complex CoT of advanced CC by extending RA .": "for sustainable IoT anomaly detection under fixed power and latency budgets.",
        "Cloud \nService": "IoT Edge\nPlatform",
        "Design and implement the necessary cryptographic technologies and key management systems for CC 's Remote Attestation appropriately.":
            "Design and implement a hardware-aware hybrid pipeline: simulator/cloud training, "
            "compressed edge inference, unified benchmarking.",
        "Root of Trust": "Training phase\n(Qiskit / PennyLane)",
        "Isolation execution environment \n(TEE)": "Compression &\nexport",
        "Remote \nAttestation": "Edge inference\n(FPGA / ARM)",
        "Memory encryption": "Qubit-depth\nconstraints",
        "・Tamper-resistant\n・Key management": "• Public datasets\n• Power metering",
        "Signature/Certificate\n·hash": "F1 / AUC /\nLatency / Energy",
        "TEE , hardware and software utilize advanced cryptographic technology.\nProvides basic (important) security features.":
            "IoT security analytics require co-design of ML models, quantum circuits, and embedded platforms.",
        "CC complexly integrated with cloud technologies , \nwe can improve and extend \nthe Chain of Trust and enable verification with Remote Attestation.":
            "Two-phase pipeline: train on simulator/cloud, compress, and infer on edge under fixed budgets.",
        "Isolation Execution Environment \n(CC)": "Hybrid Edge QML",
        "Confidential Computing": "FPGA offload",
        "By improving and expanding it , \nwe contribute to the advancement of CC.":
            "Empirical guidance on when hybrid Quantum AI outperforms classical edge AI.",
    }
    for idx in (10, 11, 12, 13, 14):
        replace_all_text(prs.slides[idx], cc_to_qai)

    # Fix slide 11 title
    replace_title_content(prs.slides[10], title="Background")

    # Slide 12 title
    replace_title_content(prs.slides[11], title="Research area")

    # Slide 13 title
    replace_title_content(prs.slides[12], title="Research objectives and challenges")

    if (IMG_DIR / "01_iot_quantum_landscape.png").exists():
        s11 = prs.slides[10]
        # Remove vendor logo pictures on background slide bottom-right if possible
        add_image(s11, IMG_DIR / "01_iot_quantum_landscape.png", Inches(0.4), Inches(2.5), Inches(4.5))

    # --- Slide 14: Related research ---
    replace_title_content(
        prs.slides[13],
        body=(
            "Prior QML for IoT anomaly detection (Chowdhury et al., AIMLSystems 2025)\n"
            "  → Competitive accuracy but substantial simulator-era overhead\n\n"
            "Projected quantum kernels for IoT (D'Amore et al., PICom 2024)\n"
            "  → Cloud-centric; limited edge deployability evidence\n\n"
            "Candidate's four-paradigm IoT survey (Babu & Mitra)\n"
            "  → PRISMA corpus n=57; unified benchmarks; TRL roadmap\n\n"
            "Embedded QML feasibility study (Dey & Raza, 2026)\n"
            "  → Hybrid architectures promising; deployment evidence still sparse"
        ),
    )
    s14 = prs.slides[13]
    for shape in iter_shapes(s14.shapes):
        if shape.has_text_frame:
            t = shape.text_frame.text
            if "Question : Can attestation" in t:
                set_text(
                    shape,
                    "RQ1: Under fixed power P and latency T, can hybrid inference "
                    "match classical edge accuracy within tolerance ε?",
                )
            if "Question : Is improvement" in t:
                set_text(
                    shape,
                    "RQ2–RQ3: Which compression strategies minimise resource use? "
                    "Does FPGA offload improve joules-per-inference vs. ARM gateways?",
                )
            if "Task 1:" in t:
                set_text(shape, "Survey gap:\nCloud-centric QML")
            if "Task 2:" in t:
                set_text(shape, "Doctoral focus:\nEdge deployability")

    # --- Slide 15 title ---
    replace_title_content(prs.slides[14], title="Approach")

    # --- Slide 16: Proposed architecture ---
    s16 = prs.slides[15]
    replace_title_content(
        s16,
        title="[ Reference ] Hybrid edge architecture (proposed)",
        body="Two-phase doctoral pipeline: Training → Compression → Edge inference → Evaluation",
    )
    clear_slide_pictures(s16)
    if (IMG_DIR / "04_bridge_goal.png").exists():
        add_image(s16, IMG_DIR / "04_bridge_goal.png", Inches(0.8), Inches(1.5), Inches(11.5))
    if (IMG_DIR / "02_promise_vs_deployability.png").exists():
        add_image(s16, IMG_DIR / "02_promise_vs_deployability.png", Inches(0.8), Inches(4.2), Inches(5.5))

    # --- Slide 17: Step A ---
    replace_title_content(
        prs.slides[16],
        title="A. Benchmark protocol & requirements (Year 1)",
        body=(
            "• PRISMA-inspired literature update; dataset selection (NSL-KDD, Edge-IIoTset, Bot-IoT)\n"
            "• Reproducible evaluation harness: accuracy, latency, algorithmic complexity, energy proxy\n"
            "• Classical edge baselines (lightweight CNN / autoencoder) under fixed budgets\n"
            "• Stakeholders: supervisor, FPGA lab, IoT testbed operators, open-source QML community"
        ),
    )

    # --- Slide 18: Step B ---
    replace_title_content(
        prs.slides[17],
        title="B. Hybrid model development & hardware integration (Years 2–3)",
        body=(
            "• VQC / quantum-kernel design with noise-aware training and qubit-depth constraints\n"
            "• Compression strategies: depth reduction, angle encoding, quantum-inspired surrogates\n"
            "• FPGA SoC (Zynq/Artix) or ARM gateway integration; power & thermal measurement\n"
            "• Head-to-head evaluation vs. classical TinyML and cloud-hybrid QML (RQ1–RQ3)"
        ),
    )

    # --- Slide 19: Timeline ---
    s19 = prs.slides[18]
    replace_title_content(
        s19,
        title="Announcement plan (48 months)",
        body=(
            "Research proceeds iteratively across benchmarking, hybrid model design, "
            "edge prototyping, and thesis synthesis. Results will be published at IEEE IoT Journal, "
            "FGCS, Applied Soft Computing, and selective conferences (IEEE IoT, QCE)."
        ),
    )
    timeline_map = {
        "Task 1\nthe TEE usage environment.": "Y1: Benchmark spec\n& classical baselines",
        "Task 2\nAdvanced CC": "Y2: Hybrid VQC\n& compression",
        "dissertation writing": "Thesis writing\n& artifact release",
        "Enrollment": "Start",
        "2024": "2026",
        "2025": "2027",
        "2026": "2028",
        "2027": "2029",
    }
    replace_all_text(s19, timeline_map)

    # --- Slide 20: Aspirations ---
    replace_title_content(
        prs.slides[19],
        title="In conclusion (Aspirations)",
        body=(
            "Regarding hybrid Quantum AI for IoT security, I will co-design quantum circuits, "
            "classical ML, and embedded platforms under real energy and latency constraints.\n\n"
            "As a security researcher, I aim to deepen expertise in deployable edge intelligence "
            "and bridge academic research with national-scale ICT practice.\n\n"
            "Ultimate goal: provide trustworthy, energy-efficient IoT security analytics and "
            "contribute to a convenient and safe intelligent society."
        ),
    )

    # --- Slide 21: Supplementary divider ---
    replace_title_content(
        prs.slides[20],
        title="Supplementary slides follow — survey highlights, four paradigms, TRL roadmap",
    )

    # --- Slides 22–29: Supplementary survey content ---
    supp = [
        (
            "[ Reference ] Survey workflow — PRISMA corpus",
            "PRISMA-inspired reduction: 229 records → 42 IoT-focused studies (+ 15 enabling refs.)\n"
            "Multi-database search: OpenAlex, Scopus, PubMed, WoS\n"
            "Tools: PyBibX, bibliometrix, VOSviewer",
        ),
        (
            "[ Reference ] Four-paradigm IoT taxonomy",
            "Conventional IoT (C:17) | AI-based IoT (AI:15)\n"
            "Quantum-enhanced (Q:4) | Quantum AI IoT (QAI:6)\n"
            "Comparison functional J(𝒲) for paradigm selection under deployability constraints",
        ),
        (
            "[ Reference ] Case studies from current survey",
            "Primary: Smart-building IoT security gateway (hybrid edge IDS)\n"
            "Secondary: Post-discharge RPM IoMT with innovative Hybrid Reference Architecture\n"
            "Policy-routed deployment paths α, β, γ for zero-trust hybrid edge",
        ),
        (
            "[ Reference ] Benchmark metrics alignment",
            "Six-panel performance simulation: latency, complexity, resource proxies, F1/AUC\n"
            "Addresses documented gap—accuracy rarely reported jointly with circuit depth and energy\n"
            "Formal propositions: qubit budgets, shot-noise, latency decomposition, memory crossover",
        ),
        (
            "[ Reference ] Strategic importance of QAI IoT",
            "Formalises hybrid necessity, Pareto improvement, and fleet-scale trade-offs\n"
            "TRL roadmap for quantum AI-based IoT fleets\n"
            "Near-sensor analytics and ultra-low-latency cyber–physical control",
        ),
        (
            "[ Reference ] Candidate publications (selected)",
            "Reversible/quantum logic: IEEE conferences; IoT–Quantum AI survey (co-PI)\n"
            "Target PhD venues: IEEE IoT Journal, FGCS, Applied Soft Computing, IEEE Access\n"
            "Conferences: IEEE IoT, QCE, edge-AI venues",
        ),
        (
            "[ Reference ] QAI IoT technology readiness (TRL)",
            "TRL 2–4 target for doctoral prototype: lab-validated hybrid edge inference\n"
            "Roadmap: simulator training → compressed export → FPGA/ARM deployment → fleet evaluation",
        ),
        (
            "[ Reference ] Publication growth in QAI IoT (2016–2026)",
            "Survey corpus shows accelerating publication growth in IoT–quantum AI intersection\n"
            "Key venues: IEEE IoT Journal, IEEE Access, Discover AI, Future Generation Computer Systems",
        ),
    ]

    supp_images = {
        0: IMG_DIR / "01_iot_quantum_landscape.png",
        1: IMG_DIR / "03_doctoral_cultivation.png",
        6: IMG_DIR / "02_promise_vs_deployability.png",
    }
    for i, (title, body) in enumerate(supp):
        slide_idx = 21 + i
        if slide_idx >= len(prs.slides):
            break
        img = supp_images.get(i)
        simplify_slide(prs.slides[slide_idx], title=title, body=body, image=img)

    fix_template_remnants(prs)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out}")
